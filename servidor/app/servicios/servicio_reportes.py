from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import io

class ServicioReportes:
    def generar_pptx(self, analisis_ia: dict) -> bytes:
        """
        Genera un archivo PPTX en memoria basado en el analisis estructurado (JSON).
        espera un array de diapositivas en analisis_ia['diapositivas']:
        [
            {
               "titulo": "Titulo de la slide",
               "puntos": ["punto 1", "punto 2"],
               "notas_orador": "Texto para las notas del orador"
            }
        ]
        """
        prs = Presentation()
        
        # Diapositiva de titulo
        slide_title_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_title_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "Informe de Análisis Estadístico"
        subtitle.text = "Generado por AgenteCraft IA"
        
        diapositivas = analisis_ia.get("diapositivas", [])
        
        # Diapositivas de contenido
        bullet_slide_layout = prs.slide_layouts[1]
        
        for d in diapositivas:
            slide = prs.slides.add_slide(bullet_slide_layout)
            shapes = slide.shapes
            
            title_shape = shapes.title
            body_shape = shapes.placeholders[1]
            
            title_shape.text = d.get("titulo", "Diapositiva")
            
            tf = body_shape.text_frame
            puntos = d.get("puntos", [])
            for i, p in enumerate(puntos):
                if i == 0:
                    tf.text = p
                else:
                    pa = tf.add_paragraph()
                    pa.text = p
            
            # Añadir notas del orador
            notas = d.get("notas_orador", "Sin notas")
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = notas

        output = io.BytesIO()
        prs.save(output)
        output.seek(0)
        return output.read()

servicio_reportes = ServicioReportes()
